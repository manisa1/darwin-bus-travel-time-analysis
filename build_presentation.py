"""
PRT564 - Assessment 2 Presentation Builder
Generates a .pptx file with all required slides and embedded visuals.
Group 10
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# PATHS / CONFIG
from config import OUTPUT_DIR, BASE_DIR

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x37, 0x5E)   # slide background / headers
MID_BLUE    = RGBColor(0x2E, 0x86, 0xAB)   # accent
ORANGE      = RGBColor(0xF1, 0x8F, 0x01)   # highlight
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xC0, 0x39, 0x2B)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────
def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(img_path), left, top, width)


def header_bar(slide, title: str, subtitle: str = ""):
    """Dark blue top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.15), DARK_BLUE)
    add_textbox(slide, title,
                Inches(0.35), Inches(0.12), Inches(10), Inches(0.65),
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(0.72), Inches(10), Inches(0.38),
                    font_size=14, color=MID_BLUE)


def slide_number(slide, n: int):
    add_textbox(slide, str(n),
                Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.3),
                font_size=11, color=DARK_GREY, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items: list[tuple[str, str]],
                 left, top, width, line_gap=Inches(0.42),
                 dot_color=ORANGE, font_size=17, text_color=DARK_GREY):
    """Render bullet items as (label, body) pairs."""
    y = top
    for label, body in items:
        # dot
        add_rect(slide, left, y + Inches(0.08), Inches(0.12), Inches(0.12), dot_color)
        # label bold
        if label:
            add_textbox(slide, label + (" — " if body else ""),
                        left + Inches(0.22), y, Inches(1.4), Inches(0.4),
                        font_size=font_size, bold=True, color=DARK_BLUE)
            add_textbox(slide, body,
                        left + Inches(1.65), y, width - Inches(1.65), Inches(0.4),
                        font_size=font_size, color=text_color)
        else:
            add_textbox(slide, body,
                        left + Inches(0.22), y, width - Inches(0.22), Inches(0.4),
                        font_size=font_size, color=text_color)
        y += line_gap


def table_slide(slide, headers, rows, left, top, col_widths, row_h=Inches(0.42)):
    """Simple table drawn as coloured rectangles."""
    x = left
    # header row
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_widths[i], row_h, DARK_BLUE)
        add_textbox(slide, h, x + Inches(0.05), top + Inches(0.05),
                    col_widths[i] - Inches(0.1), row_h - Inches(0.05),
                    font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += col_widths[i]

    for r_idx, row in enumerate(rows):
        bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
        x = left
        for i, cell in enumerate(row):
            add_rect(slide, x, top + row_h * (r_idx + 1),
                     col_widths[i], row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            c_color = GREEN if "significant" in str(cell).lower() and "not" not in str(cell).lower() else DARK_GREY
            add_textbox(slide, str(cell),
                        x + Inches(0.05),
                        top + row_h * (r_idx + 1) + Inches(0.05),
                        col_widths[i] - Inches(0.1), row_h - Inches(0.05),
                        font_size=14, color=c_color, align=PP_ALIGN.CENTER)
            x += col_widths[i]


# ════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ════════════════════════════════════════════════════════════════════════════
prs = new_prs()


# ── SLIDE 1 · Title ──────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK_BLUE)
add_rect(s, 0, Inches(2.6), W, Inches(2.5), MID_BLUE)

add_textbox(s, "Darwin Bus Network",
            Inches(1), Inches(1.0), Inches(11.3), Inches(0.9),
            font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Travel Time Prediction using Regression Modelling",
            Inches(1), Inches(1.85), Inches(11.3), Inches(0.7),
            font_size=24, color=MID_BLUE, align=PP_ALIGN.CENTER)

add_textbox(s, "PRT564 — Assessment 2  |  Group 10",
            Inches(1), Inches(2.75), Inches(11.3), Inches(0.55),
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "Research Question:",
            Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.4),
            font_size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_textbox(s,
            '"To what extent can bus travel time be predicted based on route characteristics\n'
            'such as distance travelled, number of stops, and time of day?"',
            Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.9),
            font_size=16, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "Data Sources:  NT Government GTFS Feed  |  ABS 2021 Census (SA2)  |  Geographic Reference (Darwin CBD)",
            Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
            font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 · Data Sources & Heterogeneous Integration ───────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Data Sources", "Heterogeneous data integration from 3 distinct sources")
slide_number(s, 2)

# Source cards
card_data = [
    (MID_BLUE,  "Source 1 — NT Government GTFS Feed",
     ["stop_times.txt  —  ~100k stop-event rows",
      "trips.txt  —  trip-to-route mapping",
      "stops.txt  —  GPS coordinates of each stop",
      "routes.txt  —  route names & metadata",
      "calendar.txt  —  weekday/weekend service patterns"],
     "Temporal · Spatial · Categorical"),
    (ORANGE,    "Source 2 — ABS Census 2021 (SA2)",
     ["darwin_sa2_population.csv",
      "Population counts per suburb (SA2 level)",
      "Population density per km²",
      "Covers Greater Darwin region"],
     "Demographic · External"),
    (DARK_BLUE, "Source 3 — Geographic Reference",
     ["Darwin Interchange CBD centroid",
      "Coordinates: -12.464786, 130.844340",
      "Derived from GTFS stop_id 83",
      "Used to compute Haversine distance"],
     "Spatial · Derived"),
]

for i, (color, title, bullets, dtype) in enumerate(card_data):
    cx = Inches(0.25 + i * 4.35)
    add_rect(s, cx, Inches(1.25), Inches(4.15), Inches(5.9), color)
    add_textbox(s, title, cx + Inches(0.15), Inches(1.35),
                Inches(3.85), Inches(0.55), font_size=14, bold=True, color=WHITE)
    add_rect(s, cx, Inches(1.85), Inches(4.15), Inches(0.04), WHITE)
    y = Inches(1.95)
    for b in bullets:
        add_rect(s, cx + Inches(0.18), y + Inches(0.12),
                 Inches(0.1), Inches(0.1), WHITE)
        add_textbox(s, b, cx + Inches(0.35), y,
                    Inches(3.7), Inches(0.38), font_size=13, color=WHITE)
        y += Inches(0.42)
    add_rect(s, cx, Inches(6.75), Inches(4.15), Inches(0.35), RGBColor(0,0,0))
    add_textbox(s, "Type: " + dtype, cx + Inches(0.1), Inches(6.77),
                Inches(3.9), Inches(0.3), font_size=12, bold=True, color=ORANGE)


# ── SLIDE 3 · Pipeline Flowchart ─────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Data Analysis Pipeline", "End-to-end 10-stage process")
slide_number(s, 3)

stages = [
    ("1", "Load Data",         "GTFS + ABS + Geo"),
    ("2", "Preprocess",        "Clean, parse, filter"),
    ("3", "Feature Eng.",      "Peak hours, day type"),
    ("4", "Data Integration",  "CBD dist, SA2 density"),
    ("5", "EDA",               "5 visualisations"),
    ("6", "Modelling",         "LR · DT · RF"),
    ("7", "Cross-Validation",  "5-fold + t-tests"),
    ("8", "Residual Diag.",    "Shapiro-Wilk, Q-Q"),
    ("9", "Prediction Plots",  "Actual vs Predicted"),
    ("10","Summary Output",    "CSVs + PNGs"),
]

box_w = Inches(1.15)
box_h = Inches(1.1)
gap   = Inches(0.12)
start_x = Inches(0.25)
y_top = Inches(1.4)

for i, (num, title, desc) in enumerate(stages):
    row = i // 5
    col = i % 5
    x = start_x + col * (box_w + gap)
    y = y_top + row * (box_h + Inches(1.0))
    color = MID_BLUE if row == 0 else DARK_BLUE
    add_rect(s, x, y, box_w, box_h, color)
    add_textbox(s, num, x, y + Inches(0.05), box_w, Inches(0.3),
                font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(s, title, x, y + Inches(0.3), box_w, Inches(0.35),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x, y + Inches(0.65), box_w, Inches(0.4),
                font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    # arrow (not after last in row)
    if col < 4:
        ax = x + box_w + Inches(0.01)
        add_textbox(s, "▶", ax, y + Inches(0.35), gap + Inches(0.02), Inches(0.35),
                    font_size=12, color=ORANGE, align=PP_ALIGN.CENTER)

add_textbox(s, "Heterogeneous Integration spans Stages 1–4  |  Statistical rigour ensured in Stages 6–8",
            Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
            font_size=13, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── SLIDE 4 · Preprocessing ───────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Data Preprocessing", "Every step justified by data characteristics")
slide_number(s, 4)

steps = [
    ("GTFS Time Parsing",
     "Used pd.to_timedelta (not to_datetime) because GTFS times exceed 24:00:00 for next-day services. "
     "to_datetime would silently produce wrong values."),
    ("Numeric Coercion",
     "stop_sequence and shape_dist_traveled forced to numeric; non-parseable rows dropped to prevent "
     "silent type errors propagating into aggregation."),
    ("Trip-Level Aggregation",
     "Stop-event rows aggregated per trip_id (min departure → max arrival). "
     "Required because regression target is trip travel time — an inherently trip-level quantity."),
    ("Outlier Removal (IQR)",
     "Trips outside the 1st–99th percentile of travel_time_min removed (~1% of data). "
     "Justified because Linear Regression is highly sensitive to high-leverage outliers."),
    ("Heterogeneous Join",
     "Each bus stop spatially joined to its nearest ABS SA2 centroid via Haversine distance. "
     "Stop-level features then aggregated (mean/max) to trip level for modelling."),
]

y = Inches(1.3)
for i, (title, desc) in enumerate(steps):
    color = MID_BLUE if i % 2 == 0 else DARK_BLUE
    add_rect(s, Inches(0.3), y, Inches(0.45), Inches(0.55), color)
    add_textbox(s, str(i + 1), Inches(0.3), y + Inches(0.05), Inches(0.45), Inches(0.45),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, title, Inches(0.85), y, Inches(2.5), Inches(0.55),
                font_size=14, bold=True, color=DARK_BLUE)
    add_textbox(s, desc, Inches(3.4), y, Inches(9.6), Inches(0.55),
                font_size=13, color=DARK_GREY)
    y += Inches(0.68)

add_rect(s, Inches(0.3), y + Inches(0.1), Inches(12.7), Inches(0.5), DARK_BLUE)
add_textbox(s,
            "Final dataset: 1,548 trips  |  8 features  |  3 data sources merged",
            Inches(0.3), y + Inches(0.15), Inches(12.7), Inches(0.4),
            font_size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ── SLIDE 5 · EDA — Distributions & Correlation ──────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Exploratory Data Analysis", "Distributions & Feature Correlations")
slide_number(s, 5)

add_image(s, OUTPUT_DIR / "01_distributions.png",
          Inches(0.3), Inches(1.25), Inches(6.5), Inches(3.5))
add_image(s, OUTPUT_DIR / "02_correlation_heatmap.png",
          Inches(6.9), Inches(1.2), Inches(6.2), Inches(5.7))

add_rect(s, Inches(0.3), Inches(4.85), Inches(6.5), Inches(2.3), DARK_BLUE)
insights = [
    "Travel time is roughly normally distributed (mean ~32 min)",
    "Trip distance is right-skewed — most routes under 20 km",
    "trip_distance_km has the strongest correlation with travel time (r ≈ 0.91)",
    "num_stops also strongly correlated — validates both as key predictors",
    "is_weekend_only shows low correlation — less informative feature",
]
y = Inches(4.9)
for ins in insights:
    add_rect(s, Inches(0.45), y + Inches(0.12), Inches(0.1), Inches(0.1), ORANGE)
    add_textbox(s, ins, Inches(0.62), y, Inches(6.0), Inches(0.35),
                font_size=12, color=WHITE)
    y += Inches(0.38)


# ── SLIDE 6 · EDA — Scatter & Boxplots ───────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Exploratory Data Analysis", "Spatial patterns & Time-of-day effects")
slide_number(s, 6)

add_image(s, OUTPUT_DIR / "03_scatter_distance_vs_time.png",
          Inches(0.3), Inches(1.25), Inches(6.5), Inches(3.5))
add_image(s, OUTPUT_DIR / "04_boxplots_day_type.png",
          Inches(6.9), Inches(1.25), Inches(6.2), Inches(3.5))

add_rect(s, Inches(0.3), Inches(4.85), Inches(6.5), Inches(2.3), MID_BLUE)
add_textbox(s, "Scatter: Distance vs Travel Time",
            Inches(0.45), Inches(4.9), Inches(6.2), Inches(0.4),
            font_size=14, bold=True, color=WHITE)
add_textbox(s,
            "Colour encodes mean SA2 population density. Denser suburbs (yellow) "
            "cluster at lower distances from CBD. Strong linear trend confirms "
            "trip_distance_km as the primary predictor.",
            Inches(0.45), Inches(5.3), Inches(6.2), Inches(0.9),
            font_size=13, color=WHITE)

add_rect(s, Inches(6.9), Inches(4.85), Inches(6.2), Inches(2.3), DARK_BLUE)
add_textbox(s, "Boxplots: Peak vs Weekend",
            Inches(7.05), Inches(4.9), Inches(5.9), Inches(0.4),
            font_size=14, bold=True, color=WHITE)
add_textbox(s,
            "Peak-hour trips show marginally longer times. Weekend-only services "
            "have slightly shorter median times. Both features were included but "
            "carry lower predictive weight than distance.",
            Inches(7.05), Inches(5.3), Inches(5.9), Inches(0.9),
            font_size=13, color=WHITE)


# ── SLIDE 7 · EDA — Top Routes ───────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Exploratory Data Analysis", "Route-level travel time patterns")
slide_number(s, 7)

add_image(s, OUTPUT_DIR / "05_top_routes_by_time.png",
          Inches(0.3), Inches(1.25), Inches(8.0), Inches(5.9))

add_rect(s, Inches(8.5), Inches(1.3), Inches(4.6), Inches(5.8), DARK_BLUE)
add_textbox(s, "Key Observations",
            Inches(8.65), Inches(1.4), Inches(4.3), Inches(0.45),
            font_size=16, bold=True, color=ORANGE)
insights2 = [
    "Intercity / long-haul routes dominate the top 20",
    "Significant variance within routes suggests that modelling at the trip level (not route level) is the right granularity",
    "Route identity alone is insufficient to predict travel time — distance and stops are more informative",
    "This informed the decision NOT to include route_id as a categorical predictor (too many levels, risk of overfitting)",
]
y2 = Inches(1.95)
for ins in insights2:
    add_rect(s, Inches(8.65), y2 + Inches(0.1), Inches(0.1), Inches(0.1), ORANGE)
    add_textbox(s, ins, Inches(8.82), y2, Inches(4.15), Inches(0.55),
                font_size=13, color=WHITE)
    y2 += Inches(0.72)


# ── SLIDE 8 · Regression Models ──────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Regression Models", "Three models — baseline to ensemble")
slide_number(s, 8)

models_info = [
    (MID_BLUE,  "Linear Regression",
     "Baseline model. Assumes a linear relationship between features and travel time. "
     "Interpretable coefficients allow understanding of each feature's effect. "
     "Sensitive to outliers (addressed during preprocessing).",
     "Why chosen: establishes a performance floor; coefficients interpretable for stakeholders."),
    (ORANGE,    "Decision Tree\n(max_depth=8)",
     "Captures non-linear interactions (e.g. distance × time-of-day effects). "
     "No assumption of linearity. Risk of overfitting controlled by depth limit.",
     "Why chosen: flexible, interpretable splits, feature importances available."),
    (DARK_BLUE, "Random Forest\n(200 trees, depth=12)",
     "Ensemble of decorrelated trees. Reduces variance through bagging. "
     "Most robust model — handles outliers, non-linearity, and feature interactions.",
     "Why chosen: best expected generalisation; provides reliable feature importances."),
]

for i, (color, title, desc, why) in enumerate(models_info):
    cx = Inches(0.25 + i * 4.35)
    add_rect(s, cx, Inches(1.25), Inches(4.15), Inches(5.9), color)
    add_textbox(s, title, cx + Inches(0.15), Inches(1.35),
                Inches(3.85), Inches(0.6), font_size=16, bold=True, color=WHITE)
    add_rect(s, cx, Inches(1.9), Inches(4.15), Inches(0.04), WHITE)
    add_textbox(s, desc, cx + Inches(0.15), Inches(2.0),
                Inches(3.85), Inches(2.2), font_size=13, color=WHITE)
    add_rect(s, cx, Inches(4.2), Inches(4.15), Inches(0.04), ORANGE)
    add_textbox(s, why, cx + Inches(0.15), Inches(4.3),
                Inches(3.85), Inches(1.5), font_size=12, italic=True, color=ORANGE)

add_textbox(s, "Features used (8):  trip_distance_km · num_stops · start_hour · is_peak · "
               "is_weekend_only · mean_dist_from_cbd_km · max_dist_from_cbd_km · mean_sa2_density",
            Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
            font_size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── SLIDE 9 · Model Evaluation ───────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Model Evaluation", "Hold-out test set performance (80/20 split)")
slide_number(s, 9)

headers = ["Model", "MAE (min)", "RMSE (min)", "R²", "Interpretation"]
rows = [
    ["Linear Regression", "3.40", "4.31", "0.847",
     "Explains 84.7% of variance — good baseline"],
    ["Decision Tree",     "0.87", "2.14", "0.962",
     "Strong non-linear fit, low error"],
    ["Random Forest",     "0.58", "1.71", "0.976",
     "Best model — lowest error, highest R²"],
]
col_widths = [Inches(2.4), Inches(1.4), Inches(1.4), Inches(1.0), Inches(6.8)]
table_slide(s, headers, rows,
            left=Inches(0.3), top=Inches(1.3),
            col_widths=col_widths, row_h=Inches(0.58))

# Metric justification box
add_rect(s, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.04), MID_BLUE)
add_textbox(s, "Why these metrics?",
            Inches(0.3), Inches(4.1), Inches(4.0), Inches(0.4),
            font_size=15, bold=True, color=DARK_BLUE)
metric_text = (
    "MAE — average absolute error in minutes; robust to outliers, easy for non-technical stakeholders to interpret.\n"
    "RMSE — penalises large errors more than MAE; useful when big prediction errors are especially costly.\n"
    "R² — proportion of variance explained; scale-independent, ideal for comparing models on same data."
)
add_textbox(s, metric_text,
            Inches(0.3), Inches(4.55), Inches(12.7), Inches(1.1),
            font_size=13, color=DARK_GREY)

add_image(s, OUTPUT_DIR / "07_actual_vs_predicted.png",
          Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.7))


# ── SLIDE 10 · Statistical Tests ─────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Statistical Tests", "Are performance differences significant?")
slide_number(s, 10)

# T-test table
add_textbox(s, "Paired t-tests on 5-fold CV RMSE  (H₀: equal mean RMSE between models)",
            Inches(0.3), Inches(1.25), Inches(12.7), Inches(0.4),
            font_size=15, bold=True, color=DARK_BLUE)
t_headers = ["Comparison", "Mean RMSE A", "Mean RMSE B", "t-statistic", "p-value", "Result"]
t_rows = [
    ["LR vs Decision Tree",  "4.142", "1.672", "30.50", "6.88e-06", "significant"],
    ["LR vs Random Forest",  "4.142", "1.409", "49.10", "1.03e-06", "significant"],
    ["Decision Tree vs RF",  "1.672", "1.409",  "5.23", "0.0064",   "significant"],
]
t_col_w = [Inches(3.2), Inches(1.5), Inches(1.5), Inches(1.6), Inches(1.6), Inches(3.8)]
table_slide(s, t_headers, t_rows,
            left=Inches(0.3), top=Inches(1.7), col_widths=t_col_w, row_h=Inches(0.5))

# F-test box
add_rect(s, Inches(0.3), Inches(3.4), Inches(6.0), Inches(1.6), DARK_BLUE)
add_textbox(s, "Global F-test (Linear Regression)",
            Inches(0.45), Inches(3.5), Inches(5.7), Inches(0.4),
            font_size=14, bold=True, color=ORANGE)
add_textbox(s,
            "Tests H₀: all LR coefficients = 0 (model no better than intercept-only).\n"
            "Result: F(8, ~1238) = very large,  p < 0.0001  →  model is globally significant.",
            Inches(0.45), Inches(3.95), Inches(5.7), Inches(0.9),
            font_size=13, color=WHITE)

# Shapiro box
add_rect(s, Inches(6.5), Inches(3.4), Inches(6.5), Inches(1.6), MID_BLUE)
add_textbox(s, "Shapiro-Wilk Test (Residual Normality)",
            Inches(6.65), Inches(3.5), Inches(6.2), Inches(0.4),
            font_size=14, bold=True, color=WHITE)
add_textbox(s,
            "H₀: residuals are normally distributed.\n"
            "Result: p < 0.05  →  residuals are NON-normal.\n"
            "Implication: LR predictions are valid but confidence intervals should be interpreted cautiously.",
            Inches(6.65), Inches(3.95), Inches(6.2), Inches(0.95),
            font_size=13, color=WHITE)

add_textbox(s,
            "All three models are statistically distinguishable from each other — "
            "Random Forest is the clear best performer.",
            Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.45),
            font_size=14, bold=True, italic=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_image(s, OUTPUT_DIR / "06_residual_diagnostics.png",
          Inches(0.3), Inches(5.55), Inches(12.7), Inches(1.7))


# ── SLIDE 11 · Feature Importance ─────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT_GREY)
header_bar(s, "Feature Importance & Coefficients", "What drives bus travel time?")
slide_number(s, 11)

add_image(s, OUTPUT_DIR / "08_feature_importance.png",
          Inches(0.3), Inches(1.25), Inches(8.5), Inches(3.8))

add_rect(s, Inches(8.9), Inches(1.25), Inches(4.2), Inches(3.8), DARK_BLUE)
add_textbox(s, "Linear Regression Coefficients",
            Inches(9.05), Inches(1.35), Inches(3.9), Inches(0.45),
            font_size=14, bold=True, color=ORANGE)
coefs = [
    ("trip_distance_km",       "+1.19 min/km"),
    ("num_stops",              "+0.32 min/stop"),
    ("max_dist_from_cbd_km",   "-0.16"),
    ("start_hour",             "-0.07"),
    ("is_peak",                "-0.13"),
    ("mean_sa2_density",       "+0.003"),
]
y = Inches(1.85)
for feat, val in coefs:
    add_textbox(s, feat, Inches(9.05), y, Inches(2.5), Inches(0.35), font_size=12, color=LIGHT_GREY)
    col = GREEN if "+" in val else RED
    add_textbox(s, val, Inches(11.6), y, Inches(1.4), Inches(0.35),
                font_size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)
    y += Inches(0.38)

add_rect(s, Inches(0.3), Inches(5.15), Inches(12.7), Inches(2.0), MID_BLUE)
add_textbox(s, "Interpretation for stakeholders",
            Inches(0.45), Inches(5.2), Inches(8.0), Inches(0.4),
            font_size=15, bold=True, color=WHITE)
interp = (
    "Trip distance is by far the strongest driver — every extra kilometre adds ~1.2 minutes.\n"
    "Number of stops adds ~0.3 minutes per stop — frequent-stop routes are slower.\n"
    "CBD distance has a small negative effect — outer-suburb routes travel faster per km (fewer stops, more open road).\n"
    "Peak hours and weekends have minimal impact on scheduled times in Darwin's network."
)
add_textbox(s, interp,
            Inches(0.45), Inches(5.65), Inches(12.3), Inches(1.4),
            font_size=13, color=WHITE)


# ── SLIDE 12 · Conclusion & Non-Technical Summary ────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK_BLUE)
header_bar(s, "Conclusion", "Answering the research question")
slide_number(s, 12)

add_textbox(s,
            '"To what extent can bus travel time be predicted based on route characteristics?"',
            Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.55),
            font_size=17, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Answer box
add_rect(s, Inches(0.3), Inches(1.9), Inches(12.7), Inches(1.0), MID_BLUE)
add_textbox(s,
            "Yes — very well.  Random Forest achieves R² = 0.976 (only 2.4% of variance unexplained), "
            "predicting travel time to within 1.7 minutes on average.",
            Inches(0.45), Inches(2.0), Inches(12.4), Inches(0.8),
            font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Plain-language summary
add_textbox(s, "What this means in plain language:",
            Inches(0.5), Inches(3.1), Inches(6.0), Inches(0.45),
            font_size=15, bold=True, color=ORANGE)
plain = [
    "If you know how far a bus travels and how many stops it makes, you can reliably predict how long it will take.",
    "Buses that travel further from the city centre tend to cover distance faster (fewer stops on outer routes).",
    "Peak-hour effects are small in Darwin — scheduled times don't vary much by time of day.",
    "The Random Forest model is the most accurate and can support real-time travel planning tools.",
]
y = Inches(3.6)
for p in plain:
    add_rect(s, Inches(0.5), y + Inches(0.1), Inches(0.12), Inches(0.12), ORANGE)
    add_textbox(s, p, Inches(0.72), y, Inches(5.8), Inches(0.45), font_size=13, color=WHITE)
    y += Inches(0.52)

# Limitations
add_textbox(s, "Limitations:",
            Inches(6.9), Inches(3.1), Inches(6.0), Inches(0.45),
            font_size=15, bold=True, color=ORANGE)
limits = [
    "Based on scheduled times — real delays, traffic not captured.",
    "SA2 demographic join uses nearest centroid (approximate).",
    "GTFS snapshot — does not reflect seasonal timetable changes.",
    "Random Forest is less interpretable than Linear Regression.",
]
y = Inches(3.6)
for lim in limits:
    add_rect(s, Inches(6.9), y + Inches(0.1), Inches(0.12), Inches(0.12), RED)
    add_textbox(s, lim, Inches(7.12), y, Inches(5.9), Inches(0.45), font_size=13, color=WHITE)
    y += Inches(0.52)

add_rect(s, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4), MID_BLUE)
add_textbox(s, "Group 10  |  PRT564 Data Analytics & Visualisation  |  CDU",
            Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
            font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ── SLIDE 13 · Thank You / Q&A ────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK_BLUE)
add_rect(s, 0, Inches(2.8), W, Inches(2.2), MID_BLUE)
add_textbox(s, "Thank You", Inches(1), Inches(1.6), Inches(11.3), Inches(1.0),
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Questions & Discussion",
            Inches(1), Inches(2.95), Inches(11.3), Inches(0.65),
            font_size=26, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, "Group 10  |  PRT564  |  Darwin Bus Network Travel Time Prediction",
            Inches(1), Inches(5.7), Inches(11.3), Inches(0.5),
            font_size=16, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_textbox(s, "Data: NT Government GTFS  |  ABS 2021 Census  |  Geographic Reference",
            Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
            font_size=13, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ── SAVE ─────────────────────────────────────────────────────────────────────
out_path = BASE_DIR / "PRT564_Group10_Assessment2_Presentation.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
